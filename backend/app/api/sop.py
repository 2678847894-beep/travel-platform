from fastapi import APIRouter, Depends, HTTPException, UploadFile, File, Form
from sqlalchemy.orm import Session
from typing import List, Optional
from pydantic import BaseModel
from app.core.database import get_db
from app.core.auth import get_current_user, require_admin
from app.models.user import User
from app.models.sop import SopFolder, SopDocument
from app.schemas.schemas import SopFolderOut, SopDocumentOut, SopDocumentCreate, SopDocumentUpdate
import io

router = APIRouter(prefix='/api/sop', tags=['SOP'])


class SopFolderCreate(BaseModel):
    name: str
    icon: str = "📁"
    trip_filter: str = "香港差旅"
    order_index: int = 0


@router.get('/folders', response_model=List[SopFolderOut])
def list_folders(trip_filter: str = None, db: Session = Depends(get_db)):
    q = db.query(SopFolder)
    if trip_filter and trip_filter != '全部':
        q = q.filter(SopFolder.trip_filter == trip_filter)
    return q.order_by(SopFolder.sort_order).all()


@router.post('/folders', response_model=SopFolderOut)
def create_folder(body: SopFolderCreate, db: Session = Depends(get_db), _: User = Depends(require_admin)):
    folder = SopFolder(
        name=body.name,
        icon=body.icon,
        description="",
        sort_order=body.order_index,
        trip_filter=body.trip_filter,
    )
    db.add(folder)
    db.commit()
    db.refresh(folder)
    return folder


class SopFolderUpdate(BaseModel):
    name: Optional[str] = None
    icon: Optional[str] = None


@router.put('/folders/{folder_id}', response_model=SopFolderOut)
def update_folder(folder_id: int, body: SopFolderUpdate, db: Session = Depends(get_db), _: User = Depends(require_admin)):
    folder = db.query(SopFolder).filter(SopFolder.id == folder_id).first()
    if not folder:
        raise HTTPException(status_code=404, detail='Not found')
    update_data = body.model_dump(exclude_unset=True)
    for key, val in update_data.items():
        setattr(folder, key, val)
    db.commit()
    db.refresh(folder)
    return folder


@router.delete('/folders/{folder_id}')
def delete_folder(folder_id: int, db: Session = Depends(get_db), _: User = Depends(require_admin)):
    folder = db.query(SopFolder).filter(SopFolder.id == folder_id).first()
    if not folder:
        raise HTTPException(status_code=404, detail='Not found')
    # 级联删除该文件夹下的所有文档
    db.query(SopDocument).filter(SopDocument.folder_id == folder_id).delete()
    db.delete(folder)
    db.commit()
    return {'ok': True}


class SopBulkDelete(BaseModel):
    folder_ids: List[int]


@router.post('/folders/bulk-delete')
def bulk_delete_folders(body: SopBulkDelete, db: Session = Depends(get_db), _: User = Depends(require_admin)):
    """批量删除文件夹及其文档"""
    folder_ids = body.folder_ids
    deleted_folders = 0
    deleted_docs = 0
    for folder_id in folder_ids:
        folder = db.query(SopFolder).filter(SopFolder.id == folder_id).first()
        if folder:
            doc_count = db.query(SopDocument).filter(SopDocument.folder_id == folder_id).delete()
            deleted_docs += doc_count
            db.delete(folder)
            deleted_folders += 1
    db.commit()
    return {'deleted_folders': deleted_folders, 'deleted_documents': deleted_docs}

@router.get('/documents', response_model=List[SopDocumentOut])
def list_documents(folder_id: int = None, trip_filter: str = None, db: Session = Depends(get_db)):
    q = db.query(SopDocument)
    if folder_id:
        q = q.filter(SopDocument.folder_id == folder_id)
    if trip_filter and trip_filter != 'all':
        q = q.filter(SopDocument.trip_filter == trip_filter)
    docs = q.order_by(SopDocument.sort_order).all()
    for d in docs:
        if d.folder:
            d.folder_name = d.folder.name
    return docs

@router.get('/documents/{doc_id}', response_model=SopDocumentOut)
def get_document(doc_id: int, db: Session = Depends(get_db)):
    doc = db.query(SopDocument).filter(SopDocument.id == doc_id).first()
    if not doc:
        raise HTTPException(status_code=404, detail='Not found')
    if doc.folder:
        doc.folder_name = doc.folder.name
    return doc

@router.put('/documents/{doc_id}', response_model=SopDocumentOut)
def update_document(doc_id: int, body: SopDocumentUpdate, db: Session = Depends(get_db), _: User = Depends(get_current_user)):
    doc = db.query(SopDocument).filter(SopDocument.id == doc_id).first()
    if not doc:
        raise HTTPException(status_code=404, detail='Not found')
    for key, val in body.model_dump(exclude_unset=True).items():
        setattr(doc, key, val)
    db.commit()
    db.refresh(doc)
    return doc

@router.post('/documents', response_model=SopDocumentOut)
def create_document(body: SopDocumentCreate, db: Session = Depends(get_db), _: User = Depends(require_admin)):
    doc = SopDocument(**body.model_dump())
    db.add(doc)
    db.commit()
    db.refresh(doc)
    return doc

@router.delete('/documents/{doc_id}')
def delete_document(doc_id: int, db: Session = Depends(get_db), _: User = Depends(require_admin)):
    doc = db.query(SopDocument).filter(SopDocument.id == doc_id).first()
    if not doc:
        raise HTTPException(status_code=404, detail='Not found')
    db.delete(doc)
    db.commit()
    return {'ok': True}

@router.post('/documents/{doc_id}/toggle/{step_order}')
def toggle_step(doc_id: int, step_order: int, db: Session = Depends(get_db), _: User = Depends(get_current_user)):
    doc = db.query(SopDocument).filter(SopDocument.id == doc_id).first()
    if not doc:
        raise HTTPException(status_code=404, detail='Not found')
    steps = doc.steps or []
    for s in steps:
        if s['order'] == step_order:
            s['status'] = not s.get('status', False)
    doc.steps = steps
    db.commit()
    return {'ok': True}


def _parse_file_content(filename: str, file_bytes: bytes) -> str:
    """根据文件扩展名解析文件内容为纯文本"""
    ext = filename.rsplit('.', 1)[-1].lower() if '.' in filename else ''

    if ext in ('txt', 'md', 'json', 'csv'):
        return file_bytes.decode('utf-8', errors='replace')

    if ext == 'docx':
        from docx import Document
        doc = Document(io.BytesIO(file_bytes))
        return '\n'.join(p.text for p in doc.paragraphs)

    if ext == 'pdf':
        from PyPDF2 import PdfReader
        reader = PdfReader(io.BytesIO(file_bytes))
        return '\n'.join(page.extract_text() or '' for page in reader.pages)

    if ext == 'xlsx':
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(file_bytes), read_only=True)
        lines = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            for row in ws.iter_rows(values_only=True):
                line = ' | '.join(str(cell) if cell is not None else '' for cell in row)
                if line.strip():
                    lines.append(line)
        wb.close()
        return '\n'.join(lines)

    if ext == 'pptx':
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
        return '\n'.join(texts)

    raise HTTPException(status_code=400, detail=f'不支持的文件格式: .{ext}')


@router.post('/documents/import', response_model=SopDocumentOut)
async def import_document(
    folder_id: int = Form(...),
    trip_filter: str = Form('香港差旅'),
    file: UploadFile = File(...),
    db: Session = Depends(get_db),
    _: User = Depends(require_admin),
):
    if not file.filename:
        raise HTTPException(status_code=400, detail='未选择文件')

    file_bytes = await file.read()
    content = _parse_file_content(file.filename, file_bytes)

    title = file.filename.rsplit('.', 1)[0] if '.' in file.filename else file.filename
    description = content[:500]

    doc = SopDocument(
        folder_id=folder_id,
        title=title,
        description=description,
        trip_filter=trip_filter,
        steps=[],
    )
    db.add(doc)
    db.commit()
    db.refresh(doc)
    if doc.folder:
        doc.folder_name = doc.folder.name
    return doc
